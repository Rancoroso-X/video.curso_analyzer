# video_analyzer/v4/file_processor.py
"""
Sistema de processamento multi-formato para vídeos, áudios, documentos e legendas.
Versão 4.0 - Integração completa com NASCO Analyzer
"""
from config import OUTPUT_FOLDERS
from pathlib import Path
from typing import Dict, List, Optional, Union, Tuple
import mimetypes
import streamlit as st
from dataclasses import dataclass
from enum import Enum
import json
import hashlib
import re
from datetime import datetime

# Dependências para documentos (instalar com pip)
try:
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class FileType(Enum):
    VIDEO = "video"
    AUDIO = "audio"
    DOCUMENT = "document"
    SUBTITLE = "subtitle"
    UNKNOWN = "unknown"


@dataclass
class FileInfo:
    path: Path
    type: FileType
    size: int
    duration: Optional[float] = None
    metadata: Dict = None

    def to_dict(self) -> Dict:
        """Converte para dicionário para cache."""
        return {
            'path': str(self.path),
            'type': self.type.value,
            'size': self.size,
            'duration': self.duration,
            'metadata': self.metadata or {}
        }

    @classmethod
    def from_dict(cls, data: Dict) -> 'FileInfo':
        """Cria FileInfo a partir de dicionário."""
        return cls(
            path=Path(data['path']),
            type=FileType(data['type']),
            size=data['size'],
            duration=data.get('duration'),
            metadata=data.get('metadata', {})
        )


class MultiFormatProcessor:
    """Processador de múltiplos formatos de arquivo com cache."""

    SUPPORTED_FORMATS = {
        FileType.VIDEO: ['.mp4', '.avi', '.mov', '.mkv', '.webm', '.flv', '.m4v', '.wmv', '.3gp', '.mpg', '.mpeg'],
        FileType.AUDIO: ['.mp3', '.wav', '.m4a', '.flac', '.ogg', '.aac', '.wma', '.opus'],
        FileType.DOCUMENT: ['.pdf', '.docx', '.pptx', '.rtf'],
        FileType.SUBTITLE: ['.vtt', '.ass', '.ssa', '.sub', '.sbv']
    }

    def __init__(self, base_path: Path):
        self.base_path = base_path
        self.cache_file = base_path / ".file_cache.json"
        self.detected_files: Dict[FileType, List[FileInfo]] = {
            FileType.VIDEO: [],
            FileType.AUDIO: [],
            FileType.DOCUMENT: [],
            FileType.SUBTITLE: []
        }

    def get_directory_hash(self, directory: Path) -> str:
        """Gera hash da estrutura do diretório para cache."""
        file_stats = []
        IGNORED_DIRS = set(OUTPUT_FOLDERS)
        for file_path in directory.rglob('*'):
            if file_path.is_file():
                stat = file_path.stat()
                file_stats.append(
                    f"{file_path.name}{stat.st_size}{stat.st_mtime}")

        combined = "".join(sorted(file_stats))
        return hashlib.md5(combined.encode()).hexdigest()

    def load_cache(self) -> Optional[Dict]:
        """Carrega cache se existir e for válido."""
        if not self.cache_file.exists():
            return None

        try:
            with open(self.cache_file, 'r', encoding='utf-8') as f:
                cache_data = json.load(f)

            # Verificar se cache ainda é válido
            current_hash = self.get_directory_hash(self.base_path)
            if cache_data.get('directory_hash') == current_hash:
                return cache_data
        except Exception as e:
            st.warning(f"Erro ao carregar cache: {e}")

        return None

    def save_cache(self, detected_files: Dict[FileType, List[FileInfo]]):
        """Salva resultados no cache."""
        try:
            cache_data = {
                'directory_hash': self.get_directory_hash(self.base_path),
                'timestamp': datetime.now().isoformat(),
                'detected_files': {}
            }

            for file_type, files in detected_files.items():
                cache_data['detected_files'][file_type.value] = [
                    file_info.to_dict() for file_info in files
                ]

            with open(self.cache_file, 'w', encoding='utf-8') as f:
                json.dump(cache_data, f, indent=2, ensure_ascii=False)

        except Exception as e:
            st.warning(f"Erro ao salvar cache: {e}")

    def detect_file_type(self, file_path: Path) -> FileType:
        """Detecta o tipo de arquivo baseado na extensão."""
        ext = file_path.suffix.lower()

        for file_type, extensions in self.SUPPORTED_FORMATS.items():
            if ext in extensions:
                return file_type

        return FileType.UNKNOWN

    def scan_directory(self, directory: Path, use_cache: bool = True) -> Dict[FileType, List[FileInfo]]:
        """Escaneia diretório e categoriza todos os arquivos suportados."""

        # Tentar carregar do cache primeiro
        if use_cache:
            cached_data = self.load_cache()
            if cached_data:
                try:
                    detected_files = {}
                    for file_type_str, files_data in cached_data['detected_files'].items():
                        file_type = FileType(file_type_str)
                        detected_files[file_type] = [
                            FileInfo.from_dict(file_data) for file_data in files_data
                        ]
                    return detected_files
                except Exception as e:
                    st.warning(f"Erro ao processar cache: {e}")

        # Escanear diretório
        self.detected_files = {file_type: []
                               for file_type in FileType if file_type != FileType.UNKNOWN}

        total_files = sum(1 for _ in directory.rglob('*') if _.is_file())
        progress_bar = st.progress(0)
        status_text = st.empty()

        processed = 0

        # Definir pastas ignoradas
        IGNORED_DIRS = set(OUTPUT_FOLDERS) if OUTPUT_FOLDERS else {
            'analises_ia', 'relatorios', 'logs'}

        for file_path in directory.rglob('*'):
            if file_path.is_file():
                # Verificar se o arquivo está em pasta ignorada
                if any(part in IGNORED_DIRS for part in file_path.parts):
                    continue

                processed += 1
                progress_bar.progress(processed / total_files)
                status_text.text(
                    f"Processando: {file_path.name} ({processed}/{total_files})")

                file_type = self.detect_file_type(file_path)

                if file_type != FileType.UNKNOWN:
                    file_info = FileInfo(
                        path=file_path,
                        type=file_type,
                        size=file_path.stat().st_size,
                        metadata=self._extract_metadata(file_path, file_type)
                    )

                    self.detected_files[file_type].append(file_info)
        progress_bar.empty()
        status_text.empty()

        # Salvar no cache
        if use_cache:
            self.save_cache(self.detected_files)

        return self.detected_files

    def _extract_metadata(self, file_path: Path, file_type: FileType) -> Dict:
        """Extrai metadados específicos do tipo de arquivo."""
        metadata = {
            'last_modified': datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
            'extension': file_path.suffix.lower()
        }

        try:
            if file_type == FileType.VIDEO:
                metadata.update(self._extract_video_metadata(file_path))
            elif file_type == FileType.AUDIO:
                metadata.update(self._extract_audio_metadata(file_path))
            elif file_type == FileType.DOCUMENT:
                metadata.update(self._extract_document_metadata(file_path))
            elif file_type == FileType.SUBTITLE:
                metadata.update(self._extract_subtitle_metadata(file_path))
        except Exception as e:
            metadata['error'] = str(e)

        return metadata

    def _extract_video_metadata(self, video_path: Path) -> Dict:
        """Extrai metadados de vídeo."""
        try:
            from moviepy.editor import VideoFileClip
            with VideoFileClip(str(video_path)) as clip:
                return {
                    'duration': clip.duration,
                    'fps': clip.fps if hasattr(clip, 'fps') else None,
                    'size': (clip.w, clip.h) if hasattr(clip, 'w') else None,
                    'has_audio': clip.audio is not None
                }
        except ImportError:
            # Fallback: usar ffprobe se moviepy não estiver disponível
            try:
                import subprocess
                result = subprocess.run([
                    'ffprobe', '-v', 'quiet', '-print_format', 'json',
                    '-show_format', str(video_path)
                ], capture_output=True, text=True)

                if result.returncode == 0:
                    import json
                    data = json.loads(result.stdout)
                    duration = float(data['format'].get('duration', 0))
                    return {'duration': duration, 'format': 'detected_by_ffprobe'}
            except:
                pass

            return {'error': 'Não foi possível extrair metadados do vídeo'}
        except Exception as e:
            return {'error': f"Erro ao ler vídeo: {e}"}

    def _extract_audio_metadata(self, audio_path: Path) -> Dict:
        """Extrai metadados de áudio."""
        try:
            import subprocess
            result = subprocess.run([
                'ffprobe', '-v', 'quiet', '-print_format', 'json',
                '-show_format', str(audio_path)
            ], capture_output=True, text=True)

            if result.returncode == 0:
                import json
                data = json.loads(result.stdout)
                duration = float(data['format'].get('duration', 0))
                return {
                    'duration': duration,
                    'format': audio_path.suffix,
                    'bitrate': data['format'].get('bit_rate')
                }
        except:
            pass

        return {
            'format': audio_path.suffix,
            'estimated_type': 'audio'
        }

    def _extract_document_metadata(self, doc_path: Path) -> Dict:
        """Extrai metadados de documento."""
        ext = doc_path.suffix.lower()

        try:
            if ext == '.pdf' and PDF_AVAILABLE:
                return self._pdf_metadata(doc_path)
            elif ext == '.docx' and DOCX_AVAILABLE:
                return self._docx_metadata(doc_path)
            elif ext == '.pptx' and PPTX_AVAILABLE:
                return self._pptx_metadata(doc_path)
            elif ext in ['.txt', '.md']:
                return self._text_metadata(doc_path)
            else:
                return {'pages': 1, 'type': 'document'}
        except Exception as e:
            return {'error': f"Erro ao ler documento: {e}"}

    def _extract_subtitle_metadata(self, sub_path: Path) -> Dict:
        """Extrai metadados de legenda."""
        try:
            content = sub_path.read_text(encoding='utf-8')
            lines = content.strip().split('\n')

            # Contar elementos de legenda (aproximado)
            if sub_path.suffix.lower() == '.srt':
                subtitle_count = content.count('-->')
            else:
                subtitle_count = len([line for line in lines if line.strip()])

            return {
                'subtitle_count': subtitle_count,
                'format': sub_path.suffix.lower(),
                'encoding': 'utf-8'
            }
        except:
            return {
                'format': sub_path.suffix.lower(),
                'error': 'Não foi possível ler arquivo de legenda'
            }

    def _pdf_metadata(self, pdf_path: Path) -> Dict:
        """Metadados específicos de PDF."""
        if not PDF_AVAILABLE:
            return {'error': 'PyMuPDF não instalado'}

        doc = fitz.open(str(pdf_path))
        metadata = {
            'pages': doc.page_count,
            'title': doc.metadata.get('title', ''),
            'author': doc.metadata.get('author', ''),
            'has_images': False,
            'has_text': False
        }

        # Verificar se tem imagens e texto (primeira página como amostra)
        if doc.page_count > 0:
            page = doc[0]
            image_list = page.get_images()
            text = page.get_text().strip()

            metadata['has_images'] = len(image_list) > 0
            metadata['has_text'] = len(text) > 0

        doc.close()
        return metadata

    def _docx_metadata(self, docx_path: Path) -> Dict:
        """Metadados específicos de DOCX."""
        if not DOCX_AVAILABLE:
            return {'error': 'python-docx não instalado'}

        doc = Document(str(docx_path))
        return {
            'paragraphs': len(doc.paragraphs),
            'has_tables': len(doc.tables) > 0,
            'has_images': len(doc.inline_shapes) > 0,
            'word_count': sum(len(p.text.split()) for p in doc.paragraphs)
        }

    def _pptx_metadata(self, pptx_path: Path) -> Dict:
        """Metadados específicos de PowerPoint."""
        if not PPTX_AVAILABLE:
            return {'error': 'python-pptx não instalado'}

        prs = Presentation(str(pptx_path))

        # Contar texto total
        total_text_length = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    total_text_length += len(shape.text)

        return {
            'slides': len(prs.slides),
            'layouts': len(prs.slide_layouts),
            'total_text_length': total_text_length,
            'avg_text_per_slide': total_text_length / len(prs.slides) if prs.slides else 0
        }

    def _text_metadata(self, text_path: Path) -> Dict:
        """Metadados para arquivos de texto."""
        try:
            content = text_path.read_text(encoding='utf-8')
            lines = content.split('\n')
            words = content.split()

            return {
                'lines': len(lines),
                'words': len(words),
                'characters': len(content),
                'encoding': 'utf-8'
            }
        except UnicodeDecodeError:
            return {
                'encoding': 'unknown',
                'error': 'Erro de codificação'
            }


class DocumentExtractor:
    """Extrator de texto de documentos com suporte avançado."""

    def __init__(self):
        self.supported_formats = []

        if PDF_AVAILABLE:
            self.supported_formats.extend(['.pdf'])
        if DOCX_AVAILABLE:
            self.supported_formats.extend(['.docx'])
        if PPTX_AVAILABLE:
            self.supported_formats.extend(['.pptx'])

        self.supported_formats.extend(['.txt', '.md', '.rtf'])

    def can_extract(self, file_path: Path) -> bool:
        """Verifica se pode extrair texto do arquivo."""
        return file_path.suffix.lower() in self.supported_formats

    def extract_text(self, file_path: Path) -> str:
        """Extrai texto de qualquer documento suportado."""
        ext = file_path.suffix.lower()

        extractors = {
            '.pdf': self._extract_from_pdf,
            '.docx': self._extract_from_docx,
            '.pptx': self._extract_from_pptx,
            '.txt': self._extract_from_txt,
            '.md': self._extract_from_txt,
            '.rtf': self._extract_from_txt
        }

        extractor = extractors.get(ext)
        if extractor:
            return extractor(file_path)
        else:
            raise ValueError(f"Formato não suportado: {ext}")

    def _extract_from_txt(self, txt_path: Path) -> str:
        """Le texto puro de arquivos .txt"""
        return txt_path.read_text(encoding='utf-8', errors='ignore')

    def _extract_from_pdf(self, pdf_path: Path) -> str:
        """Extrai texto de PDF usando PyMuPDF."""
        if not PDF_AVAILABLE:
            raise ImportError(
                "PyMuPDF não instalado. Execute: pip install PyMuPDF")

        doc = fitz.open(str(pdf_path))
        text_parts = []

        for page_num in range(doc.page_count):
            page = doc[page_num]
            page_text = page.get_text()

            if page_text.strip():
                text_parts.append(f"\n--- Página {page_num + 1} ---\n")
                text_parts.append(page_text)

        doc.close()
        return "\n".join(text_parts).strip()

    def _extract_from_docx(self, docx_path: Path) -> str:
        """Extrai texto de DOCX."""
        if not DOCX_AVAILABLE:
            raise ImportError(
                "python-docx não instalado. Execute: pip install python-docx")

        doc = Document(str(docx_path))
        text_parts = []

        # Extrair parágrafos
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        # Extrair tabelas
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    table_text.append(" | ".join(row_text))

            if table_text:
                text_parts.append("\n--- Tabela ---\n")
                text_parts.extend(table_text)

        return "\n".join(text_parts)

    def _extract_from_pptx(self, pptx_path: Path) -> str:
        """Extrai texto de PowerPoint."""
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx não instalado. Execute: pip install python-pptx")

        prs = Presentation(str(pptx_path))
        text_parts = []

        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n--- Slide {i} ---")

            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                text_parts.extend(slide_text)
            else:
                text_parts.append("[Slide sem texto detectável]")

        return "\n".join(text_parts)

   # Função corrigida para render_file_upload_zone() no file_processor.py


def render_file_upload_zone():
    """Renderiza zona de upload com drag & drop melhorada - CORRIGIDO."""

    st.markdown("""
    <div class="file-upload-container">
        <div class="drag-drop-zone" id="drag-drop-zone">
            <div class="drag-drop-content">
                <div class="upload-icon">📁</div>
                <h3>Arraste e solte arquivos aqui</h3>
                <p>Suporte total para vídeos, áudios, PDFs, Word, PowerPoint e legendas</p>
                <div class="supported-formats">
                    <span class="format-badge">MP4</span>
                    <span class="format-badge">PDF</span>
                    <span class="format-badge">DOCX</span>
                    <span class="format-badge">MP3</span>
                    <span class="format-badge">PPTX</span>
                    <span class="format-badge">SRT</span>
                    <span class="format-badge">+25 formatos</span>
                </div>
                <div class="upload-note">
                    <p>💡 <strong>Nota:</strong> Esta é uma demonstração visual da zona de upload v4.0 ULTIMATE</p>
                    <p>Para carregar arquivos, use o campo "Caminho da pasta" na barra lateral ←</p>
                </div>
            </div>
        </div>
    </div>
    
    <style>
    .file-upload-container {
        margin: 2rem 0;
    }
    
    .drag-drop-zone {
        border: 3px dashed #007bff;
        border-radius: 16px;
        padding: 3rem 2rem;
        text-align: center;
        background: linear-gradient(135deg, rgba(0, 123, 255, 0.05) 0%, rgba(0, 123, 255, 0.1) 100%);
        transition: all 0.3s ease;
        cursor: pointer;
        position: relative;
    }
    
    .drag-drop-zone:hover {
        border-color: #0056b3;
        background: linear-gradient(135deg, rgba(0, 123, 255, 0.1) 0%, rgba(0, 123, 255, 0.2) 100%);
        transform: translateY(-2px);
        box-shadow: 0 8px 25px rgba(0, 123, 255, 0.15);
    }
    
    .upload-icon {
        font-size: 4rem;
        margin-bottom: 1rem;
        opacity: 0.7;
        animation: pulse 2s infinite;
    }
    
    @keyframes pulse {
        0% { transform: scale(1); }
        50% { transform: scale(1.05); }
        100% { transform: scale(1); }
    }
    
    .drag-drop-content h3 {
        color: #007bff;
        margin: 1rem 0 0.5rem 0;
        font-weight: 600;
    }
    
    .drag-drop-content p {
        color: #666;
        margin-bottom: 1.5rem;
        font-size: 0.9rem;
    }
    
    .supported-formats {
        display: flex;
        justify-content: center;
        gap: 0.5rem;
        margin-bottom: 2rem;
        flex-wrap: wrap;
    }
    
    .format-badge {
        background: rgba(0, 123, 255, 0.1);
        color: #007bff;
        padding: 0.3rem 0.8rem;
        border-radius: 20px;
        font-size: 0.8rem;
        font-weight: 500;
        transition: all 0.3s ease;
    }
    
    .format-badge:hover {
        background: rgba(0, 123, 255, 0.2);
        transform: scale(1.05);
    }
    
    .upload-note {
        background: rgba(255, 193, 7, 0.1);
        border: 1px solid rgba(255, 193, 7, 0.3);
        border-radius: 8px;
        padding: 1rem;
        margin-top: 2rem;
    }
    
    .upload-note p {
        margin: 0.5rem 0;
        font-size: 0.85rem;
        color: #856404;
    }
    
    /* Tema escuro */
    [data-theme="dark"] .drag-drop-zone {
        border-color: #90caf9;
        background: linear-gradient(135deg, rgba(144, 202, 249, 0.1) 0%, rgba(144, 202, 249, 0.15) 100%);
    }
    
    [data-theme="dark"] .drag-drop-content h3 {
        color: #90caf9;
    }
    
    [data-theme="dark"] .drag-drop-content p {
        color: #b0b0b0;
    }
    
    [data-theme="dark"] .format-badge {
        background: rgba(144, 202, 249, 0.2);
        color: #90caf9;
    }
    
    [data-theme="dark"] .upload-note {
        background: rgba(255, 193, 7, 0.15);
        border-color: rgba(255, 193, 7, 0.4);
    }
    
    [data-theme="dark"] .upload-note p {
        color: #ffc107;
    }
    </style>
    """, unsafe_allow_html=True)


def render_detected_files_summary(detected_files: Dict[FileType, List[FileInfo]]):
    """Renderiza resumo avançado dos arquivos detectados."""

    if not any(detected_files.values()):
        st.info("👆 Nenhum arquivo detectado. Faça upload ou selecione uma pasta.")
        return

    st.markdown("## 📊 Arquivos Detectados - Análise Completa")

    # Cards de resumo por tipo
    col1, col2, col3, col4 = st.columns(4)

    type_info = {
        FileType.VIDEO: ("🎥", "Vídeos", col1, "#e74c3c"),
        FileType.AUDIO: ("🎵", "Áudios", col2, "#f39c12"),
        FileType.DOCUMENT: ("📄", "Documentos", col3, "#3498db"),
        FileType.SUBTITLE: ("📝", "Legendas", col4, "#2ecc71")
    }

    for file_type, (icon, label, col, color) in type_info.items():
        count = len(detected_files[file_type])
        total_size = sum(f.size for f in detected_files[file_type])

        with col:
            st.markdown(f"""
            <div style="
                background: linear-gradient(135deg, {color}15 0%, {color}25 100%);
                padding: 1.5rem;
                border-radius: 12px;
                border-left: 4px solid {color};
                text-align: center;
                margin-bottom: 1rem;
                transition: transform 0.3s ease;
            " onmouseover="this.style.transform='translateY(-2px)'" onmouseout="this.style.transform='translateY(0)'">
                <div style="font-size: 2.5rem; margin-bottom: 0.5rem;">{icon}</div>
                <div style="font-size: 1.8rem; font-weight: bold; color: {color};">{count}</div>
                <div style="font-size: 0.9rem; color: #666; margin: 0.3rem 0;">{label}</div>
                <div style="font-size: 0.8rem; color: #888;">{format_file_size(total_size)}</div>
            </div>
            """, unsafe_allow_html=True)

    # Análise detalhada por tipo
    for file_type, files in detected_files.items():
        if files:
            icon, label, _, color = type_info[file_type]

            with st.expander(f"{icon} {label} - {len(files)} arquivo(s) | {format_file_size(sum(f.size for f in files))}"):

                # Estatísticas rápidas
                if file_type == FileType.VIDEO:
                    total_duration = sum(f.metadata.get('duration', 0)
                                         for f in files if f.metadata.get('duration'))
                    st.info(
                        f"⏱️ Duração total dos vídeos: {format_duration(total_duration)}")

                elif file_type == FileType.DOCUMENT:
                    total_pages = sum(f.metadata.get('pages', 0)
                                      for f in files if f.metadata.get('pages'))
                    st.info(f"📄 Total de páginas: {total_pages}")

                # Lista detalhada de arquivos
                # Mostrar até 15 arquivos
                for i, file_info in enumerate(files[:15]):
                    col_name, col_size, col_meta, col_actions = st.columns([
                                                                           3, 1, 2, 1])

                    with col_name:
                        st.write(f"📁 **{file_info.path.name}**")
                        st.caption(f"📂 {file_info.path.parent.name}")

                    with col_size:
                        st.write(format_file_size(file_info.size))

                    with col_meta:
                        if file_info.metadata:
                            meta_str = []

                            if 'duration' in file_info.metadata:
                                meta_str.append(
                                    f"⏱️ {format_duration(file_info.metadata['duration'])}")
                            if 'pages' in file_info.metadata:
                                meta_str.append(
                                    f"📄 {file_info.metadata['pages']} pgs")
                            if 'slides' in file_info.metadata:
                                meta_str.append(
                                    f"📊 {file_info.metadata['slides']} slides")
                            if 'words' in file_info.metadata:
                                meta_str.append(
                                    f"📝 {file_info.metadata['words']} palavras")

                            st.write(" • ".join(meta_str))

                    with col_actions:
                        if st.button("🔍", key=f"inspect_{i}_{file_type.value}", help="Inspecionar arquivo"):
                            st.info(
                                f"Arquivo: {file_info.path}\nTipo: {file_info.type.value}\nMetadados: {file_info.metadata}")

                if len(files) > 15:
                    st.caption(f"... e mais {len(files) - 15} arquivo(s)")


def format_file_size(size_bytes: int) -> str:
    """Formata tamanho do arquivo em formato legível."""
    if size_bytes == 0:
        return "0 B"

    size_names = ["B", "KB", "MB", "GB", "TB"]
    i = 0
    while size_bytes >= 1024 and i < len(size_names) - 1:
        size_bytes /= 1024.0
        i += 1

    return f"{size_bytes:.1f} {size_names[i]}"


def format_duration(seconds: float) -> str:
    """Formata duração em formato legível."""
    if not seconds or seconds <= 0:
        return "N/A"

    hours = int(seconds // 3600)
    minutes = int((seconds % 3600) // 60)
    seconds = int(seconds % 60)

    if hours > 0:
        return f"{hours}h {minutes}m {seconds}s"
    elif minutes > 0:
        return f"{minutes}m {seconds}s"
    else:
        return f"{seconds}s"


def check_dependencies() -> Dict[str, bool]:
    """Verifica quais dependências estão instaladas."""
    return {
        'PyMuPDF (PDF)': PDF_AVAILABLE,
        'python-docx (Word)': DOCX_AVAILABLE,
        'python-pptx (PowerPoint)': PPTX_AVAILABLE
    }


def render_dependency_status():
    """Renderiza status das dependências."""
    deps = check_dependencies()

    st.markdown("### 📦 Status das Dependências")

    cols = st.columns(len(deps))
    for i, (dep_name, available) in enumerate(deps.items()):
        with cols[i]:
            if available:
                st.success(f"✅ {dep_name}")
            else:
                st.error(f"❌ {dep_name}")
                st.caption("Execute: pip install [dependência]")

# ===================================================================
# INTEGRAÇÃO COM SISTEMA ATUAL
# ===================================================================


def mapear_modulos_multiformat(caminho: str, detected_files: Dict[FileType, List['FileInfo']] = None) -> dict:
    """Versão expandida do mapear_modulos para múltiplos formatos."""
    base_path = Path(caminho)
    extractor = DocumentExtractor()

    # Se detected_files não foi fornecido, escanear diretório
    if detected_files is None:
        processor = MultiFormatProcessor(base_path)
        detected_files = processor.scan_directory(base_path)

    modulos_mapeados = {}

    # Processar vídeos
    # Usar .get para segurança
    for video_file_info in detected_files.get(FileType.VIDEO, []):
        module_path = video_file_info.path.parent
        module_name = module_path.name if module_path != base_path else "modulo_raiz"

        if module_name not in modulos_mapeados:
            modulos_mapeados[module_name] = []

        stem = video_file_info.path.stem
        related_files = find_related_files(
            video_file_info.path, detected_files)

        aula_info = {
            'stem': stem,
            'video_path': str(video_file_info.path),
            'txt_path': related_files.get('transcription'),
            'srt_path': related_files.get('subtitle'),
            'metadata': video_file_info.metadata,
            'type': 'video'
        }
        modulos_mapeados[module_name].append(aula_info)

    # Processar áudios independentes
    # Usar .get para segurança
    for audio_file_info in detected_files.get(FileType.AUDIO, []):
        is_extracted_audio = False
        for video_file_in_module in detected_files.get(FileType.VIDEO, []):
            if video_file_in_module.path.parent == audio_file_info.path.parent and video_file_in_module.path.stem == audio_file_info.path.stem:
                is_extracted_audio = True
                break

        if not is_extracted_audio:
            module_path = audio_file_info.path.parent
            module_name = f"{module_path.name}_audios" if module_path != base_path else "audios_raiz"

            if module_name not in modulos_mapeados:
                modulos_mapeados[module_name] = []

            aula_info = {
                'stem': audio_file_info.path.stem,
                'video_path': None,
                'audio_path': str(audio_file_info.path),
                'txt_path': None,
                'metadata': audio_file_info.metadata,
                'type': 'audio'
            }
            modulos_mapeados[module_name].append(aula_info)

    # Processar documentos
    # Usar .get para segurança
    for doc_file_info in detected_files.get(FileType.DOCUMENT, []):
        module_path = doc_file_info.path.parent
        module_name = f"{module_path.name}_docs" if module_path != base_path else "documentos_raiz"

        if module_name not in modulos_mapeados:
            modulos_mapeados[module_name] = []

        txt_path_str = None
        if extractor.can_extract(doc_file_info.path):
            try:
                doc_text = extractor.extract_text(doc_file_info.path)
                temp_txt_path = doc_file_info.path.with_suffix('.txt')
                temp_txt_path.write_text(doc_text, encoding='utf-8')
                txt_path_str = str(temp_txt_path)
            except Exception as e:
                st.warning(
                    f"Erro ao processar documento {doc_file_info.path.name}: {e}")

        aula_info = {
            'stem': doc_file_info.path.stem,
            'video_path': None,
            'txt_path': txt_path_str,
            'doc_path': str(doc_file_info.path),
            'metadata': doc_file_info.metadata,
            'type': 'document'
        }
        modulos_mapeados[module_name].append(aula_info)

    return modulos_mapeados


def find_related_files(main_file_path: Path, detected_files: Dict[FileType, List['FileInfo']]) -> Dict:
    """
    Encontra arquivos relacionados (legendas, transcrições) na lista de arquivos detectados.
    Recebe main_file_path (o Path do vídeo principal) e a lista global detected_files.
    """
    main_stem = main_file_path.stem
    main_parent = main_file_path.parent
    related = {}

    # Buscar transcrição existente (.txt)
    txt_path = main_parent / f"{main_stem}.txt"
    if txt_path.exists():
        related['transcription'] = str(txt_path)

    # Buscar legendas na lista de SUBTITLE Files
    # Usar .get para segurança
    for subtitle_file_info in detected_files.get(FileType.SUBTITLE, []):
        if subtitle_file_info.path.stem == main_stem and subtitle_file_info.path.parent == main_parent:
            related['subtitle'] = str(subtitle_file_info.path)
            if 'transcription' not in related:
                try:
                    subtitle_text = convert_subtitle_to_text(
                        subtitle_file_info.path)
                    if subtitle_text:
                        temp_txt_path = main_parent / f"{main_stem}.txt"
                        temp_txt_path.write_text(
                            subtitle_text, encoding='utf-8')
                        related['transcription'] = str(temp_txt_path)
                except Exception as e:
                    st.warning(
                        f"Erro ao converter legenda {subtitle_file_info.path.name}: {e}")
            break

    return related


def convert_subtitle_to_text(subtitle_path: Path) -> str:
    """Converte arquivo de legenda para texto simples."""
    try:
        content = subtitle_path.read_text(encoding='utf-8')

        if subtitle_path.suffix.lower() == '.srt':
            import re
            text = re.sub(
                r'\d+\n\d{2}:\d{2}:\d{2},\d{3} --> \d{2}:\d{2}:\d{2},\d{3}\n', '', content)
            text = re.sub(r'\n\s*\n', '\n', text)
            return text.strip()

        elif subtitle_path.suffix.lower() == '.vtt':
            lines = content.split('\n')
            text_lines = []
            for line in lines:
                line = line.strip()
                if line and not line.startswith('WEBVTT') and '-->' not in line and not line.isdigit():
                    text_lines.append(line)
            return '\n'.join(text_lines)

        else:
            lines = content.split('\n')
            text_lines = []
            for line in lines:
                line = line.strip()
                if line and not re.match(r'^\d+$', line) and '-->' not in line:
                    text_lines.append(line)
            return '\n'.join(text_lines)

    except Exception as e:
        st.error(f"Erro ao converter legenda {subtitle_path.name}: {e}")
        return ""
